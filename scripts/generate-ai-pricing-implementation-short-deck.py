#!/usr/bin/env python3
"""Generate short AI pricing implementation deck (.pptx)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs/presentations/AI_Pricing_Implementation_Short_Deck.pptx"

NAVY = RGBColor(15, 23, 42)
TEAL = RGBColor(13, 148, 136)
BLUE = RGBColor(56, 189, 248)
SLATE = RGBColor(71, 85, 105)
LIGHT = RGBColor(241, 245, 249)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(22, 163, 74)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_box(slide, left, top, width, height, text, size=18, bold=False, color=SLATE, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return shape


def add_bullets(slide, left, top, width, height, items, size=16, color=SLATE):
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.word_wrap = True
    for index, item in enumerate(items):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
        p.level = 0
        p.bullet = True
    return shape


def add_title(slide, title, subtitle=None, dark=False):
    add_box(
        slide,
        Inches(0.6),
        Inches(0.45),
        Inches(12),
        Inches(0.7),
        title,
        size=30,
        bold=True,
        color=WHITE if dark else NAVY,
    )
    line = slide.shapes.add_shape(1, Inches(0.6), Inches(1.13), Inches(1.15), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()
    if subtitle:
        add_box(
            slide,
            Inches(0.6),
            Inches(1.3),
            Inches(12),
            Inches(0.5),
            subtitle,
            size=15,
            color=RGBColor(203, 213, 225) if dark else SLATE,
        )


def add_step(slide, number, title, detail, x, y):
    circle = slide.shapes.add_shape(1, x, y, Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = TEAL
    circle.line.fill.background()
    add_box(slide, x, y + Inches(0.08), Inches(0.55), Inches(0.35), number, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_box(slide, x - Inches(0.15), y + Inches(0.75), Inches(2.1), Inches(0.35), title, size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_box(slide, x - Inches(0.25), y + Inches(1.18), Inches(2.3), Inches(0.85), detail, size=12, color=SLATE, align=PP_ALIGN.CENTER)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    add_box(slide, Inches(0.7), Inches(1.6), Inches(11.8), Inches(1), "AI-Guided Smart Pricing", size=42, bold=True, color=WHITE)
    add_box(slide, Inches(0.7), Inches(2.5), Inches(11.5), Inches(0.8), "Implementation process, workflow, and business benefit", size=22, color=RGBColor(203, 213, 225))
    add_box(slide, Inches(0.7), Inches(5.8), Inches(11.5), Inches(0.4), "RipX / Echologyx  •  Short Client Approval Deck", size=14, color=TEAL)

    # 2. Goal
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "Goal of the New System", "Move from manual price testing to AI-guided profit optimization.")
    add_bullets(
        slide,
        Inches(0.8),
        Inches(2),
        Inches(11.5),
        Inches(4),
        [
            "AI recommends which products to test and what price bands to try.",
            "Merchants approve before anything goes live.",
            "RipX runs a real checkout-safe price experiment.",
            "The system learns from results and recommends the next best action.",
            "Primary metric is profit per visitor, not conversion alone.",
        ],
        size=18,
    )

    # 3. Core process
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "How It Will Work", "A simple repeatable flow from data to winner.")
    steps = [
        ("1", "Analyze", "Orders, product data, margin, and test history"),
        ("2", "Recommend", "AI ranks opportunities and explains why"),
        ("3", "Approve", "Merchant reviews prices and guardrails"),
        ("4", "Test", "RipX launches controlled A/B price test"),
        ("5", "Apply", "Winner is applied or next round is suggested"),
    ]
    for idx, step in enumerate(steps):
        add_step(slide, *step, x=Inches(0.75 + idx * 2.45), y=Inches(2.35))

    # 4. Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "Implementation Architecture", "AI advises; RipX validates and executes.")
    columns = [
        ("Data Layer", "Shopify orders\nProduct prices\nCOGS / margin\nRipX analytics"),
        ("AI Layer", "Structured OpenAI output\nOpportunity ranking\nRecommendation summary\nReasoning"),
        ("Safety Layer", "Margin floors\nMax price change\nSample size rules\nApproval gates"),
        ("Execution Layer", "RipX price test\nCheckout alignment\nAnalytics\nWinner rollout"),
    ]
    for idx, (title, body) in enumerate(columns):
        x = Inches(0.65 + idx * 3.1)
        rect = slide.shapes.add_shape(1, x, Inches(2.0), Inches(2.75), Inches(3.4))
        rect.fill.solid()
        rect.fill.fore_color.rgb = WHITE
        rect.line.color.rgb = TEAL
        add_box(slide, x + Inches(0.2), Inches(2.25), Inches(2.35), Inches(0.4), title, size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_box(slide, x + Inches(0.25), Inches(3.0), Inches(2.25), Inches(1.8), body, size=13, color=SLATE, align=PP_ALIGN.CENTER)

    # 5. AI API usage
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "How We Use the AI API", "OpenAI provides structured recommendations, not uncontrolled price changes.")
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.85),
        Inches(6),
        Inches(4.5),
        [
            "Backend calculates allowed candidate prices first.",
            "OpenAI receives compact metrics and valid candidate prices.",
            "OpenAI returns strict JSON: ranked plan, reason, risk, price arms.",
            "RipX validates every field before creating a test.",
        ],
        size=16,
    )
    add_box(slide, Inches(7.25), Inches(2.2), Inches(5), Inches(2.5), "Key Rule", size=24, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    add_box(
        slide,
        Inches(7.25),
        Inches(3.0),
        Inches(5),
        Inches(1.4),
        "AI can recommend.\nOnly RipX can validate, launch, and apply.",
        size=20,
        bold=True,
        color=NAVY,
        align=PP_ALIGN.CENTER,
    )

    # 6. Optimization loop
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "Continuous Optimization Loop", "The system improves through controlled rounds.")
    add_bullets(
        slide,
        Inches(0.8),
        Inches(1.85),
        Inches(11.5),
        Inches(4.5),
        [
            "Round 1: test current price vs small increase/decrease.",
            "If a price wins, AI recommends the next test around that winner.",
            "Each round narrows closer to the best profit point.",
            "Later, bandit mode can shift traffic toward stronger variants automatically.",
            "Fully automatic next-round testing is optional and requires merchant opt-in.",
        ],
        size=17,
    )

    # 7. Guardrails
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "Safety and Control", "Pricing is sensitive, so guardrails are built into the process.")
    add_bullets(
        slide,
        Inches(0.8),
        Inches(1.85),
        Inches(11.5),
        Inches(4.5),
        [
            "Minimum margin and floor price protection.",
            "Maximum price increase/decrease per test.",
            "No catalog price update without approval in v1.",
            "Auto-pause if conversion or profit drops too sharply.",
            "Audit log for AI recommendation, merchant decision, and final action.",
        ],
        size=17,
    )

    # 8. Implementation phases
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "Implementation Plan", "Build in safe phases.")
    phases = [
        ("Phase 0", "Data foundation", "Order sync, SKU metrics, opportunities"),
        ("Phase 1", "AI MVP", "OpenAI recommendations + one-click test launch"),
        ("Phase 2", "Learning", "Next-round recommendations + profit stats"),
        ("Phase 3", "Automation", "Optional bandit mode + approval-gated rollout"),
    ]
    for idx, (phase, title, body) in enumerate(phases):
        y = Inches(1.8 + idx * 1.15)
        add_box(slide, Inches(0.8), y, Inches(1.6), Inches(0.35), phase, size=16, bold=True, color=TEAL)
        add_box(slide, Inches(2.5), y, Inches(3.0), Inches(0.35), title, size=16, bold=True, color=NAVY)
        add_box(slide, Inches(5.4), y, Inches(6.8), Inches(0.35), body, size=15, color=SLATE)

    # 9. Benefits
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "Business Benefits", "Why this is valuable for clients.")
    add_bullets(
        slide,
        Inches(0.8),
        Inches(1.85),
        Inches(11.5),
        Inches(4.5),
        [
            "Faster pricing decisions with less manual setup.",
            "Better margin control by optimizing profit per visitor.",
            "Confidence from real customer experiments, not guesswork.",
            "Reusable learning loop across products and collections.",
            "Strong differentiation: AI-guided pricing backed by Echologyx experimentation expertise.",
        ],
        size=17,
    )

    # 10. Recommendation
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    add_title(slide, "Recommendation", "Start with a pilot MVP.", dark=True)
    add_bullets(
        slide,
        Inches(0.8),
        Inches(1.9),
        Inches(11.5),
        Inches(4.5),
        [
            "Approve OpenAI-powered Smart Pricing as the new price optimization direction.",
            "Build Phase 0 + Phase 1 first: data foundation and AI-recommended test launch.",
            "Pilot with 1–2 stores before enabling wider automation.",
            "Keep manual price tests as Advanced mode.",
        ],
        size=19,
        color=WHITE,
    )
    add_box(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.5), "AI recommends. RipX validates. Merchants approve. Results improve over time.", size=16, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Wrote {OUTPUT}")


if __name__ == "__main__":
    build()
