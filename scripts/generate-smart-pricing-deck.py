#!/usr/bin/env python3
"""Generate client-facing Smart Pricing strategy deck (.pptx)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs/presentations/AI_Smart_Pricing_Client_Deck.pptx"

# Echologyx-adjacent palette
NAVY = RGBColor(0x0F, 0x17, 0x2A)
TEAL = RGBColor(0x0D, 0x94, 0x88)
SLATE = RGBColor(0x47, 0x55, 0x69)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
ACCENT = RGBColor(0x38, 0xBD, 0xF8)


def set_slide_bg(slide, color=NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_light_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT)
    return slide


def add_dark_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    return slide


def textbox(slide, left, top, width, height, text, size=18, bold=False, color=SLATE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def bullets(slide, left, top, width, height, items, size=16, color=SLATE, spacing=Pt(8)):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = spacing
        p.bullet = True
    return box


def slide_title(slide, title, subtitle=None, dark=False):
    title_color = WHITE if dark else NAVY
    sub_color = RGBColor(0xCB, 0xD5, 0xE1) if dark else SLATE
    textbox(slide, Inches(0.7), Inches(0.45), Inches(11.8), Inches(0.8), title, size=32, bold=True, color=title_color)
    if subtitle:
        textbox(slide, Inches(0.7), Inches(1.15), Inches(11.5), Inches(0.6), subtitle, size=16, color=sub_color)


def accent_bar(slide, dark=False):
    bar = slide.shapes.add_shape(1, Inches(0.7), Inches(1.05 if not dark else 1.0), Inches(1.2), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 Title
    slide = add_dark_slide(prs)
    accent_bar(slide, dark=True)
    textbox(slide, Inches(0.7), Inches(1.6), Inches(11.5), Inches(1.2), "AI-Guided Smart Pricing", size=44, bold=True, color=WHITE)
    textbox(
        slide,
        Inches(0.7),
        Inches(2.5),
        Inches(11),
        Inches(1),
        "From manual price tests to profit-first, AI-recommended optimization",
        size=22,
        color=RGBColor(0xCB, 0xD5, 0xE1),
    )
    textbox(slide, Inches(0.7), Inches(5.8), Inches(8), Inches(0.5), "RipX / Echologyx  •  Client Strategy Deck  •  2026", size=14, color=TEAL)

    # 2 Executive summary
    slide = add_light_slide(prs)
    slide_title(slide, "Executive Summary")
    accent_bar(slide)
    bullets(
        slide,
        Inches(0.7),
        Inches(1.6),
        Inches(11.8),
        Inches(5),
        [
            "Reposition RipX: AI recommends what to test; merchants stay in control.",
            "Goal: maximize profit per visitor — not conversion alone.",
            "Method: controlled A/B experiments with guardrails — not blind auto-repricing.",
            "Foundation already exists: price tests, checkout alignment, analytics, Shopify API.",
            "New build: pricing intelligence layer (order data + recommendations + Smart Pricing UI).",
            "No new external pricing platform required for MVP.",
        ],
        size=18,
    )

    # 3 Problem
    slide = add_light_slide(prs)
    slide_title(slide, "The Problem Today")
    accent_bar(slide)
    bullets(
        slide,
        Inches(0.7),
        Inches(1.6),
        Inches(5.5),
        Inches(5),
        [
            "Merchants guess prices or copy competitors.",
            "Manual A/B setup is slow and requires expertise.",
            "Conversion-only metrics push toward discounts.",
            "Pricing decisions lack a repeatable learning system.",
        ],
        size=17,
    )
    textbox(slide, Inches(6.8), Inches(2.2), Inches(5.5), Inches(3.5), "Opportunity", size=22, bold=True, color=TEAL)
    bullets(
        slide,
        Inches(6.8),
        Inches(2.9),
        Inches(5.5),
        Inches(3.5),
        [
            "Encode Echologyx experimentation expertise in software.",
            "Lead with AI-first UX; manual tests become Advanced mode.",
            "Differentiate on profit focus + checkout-accurate pricing.",
        ],
        size=16,
        color=NAVY,
    )

    # 4 Vision
    slide = add_dark_slide(prs)
    slide_title(slide, "Our Vision", dark=True)
    accent_bar(slide, dark=True)
    textbox(
        slide,
        Inches(0.7),
        Inches(2.0),
        Inches(11.5),
        Inches(2.5),
        '"AI finds your best price.\nYou stay in control."',
        size=36,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    textbox(
        slide,
        Inches(1.2),
        Inches(4.5),
        Inches(10.5),
        Inches(1.5),
        "Primary flow: Smart Pricing  →  Review AI plan  →  Launch test  →  Learn  →  Apply winner",
        size=18,
        color=RGBColor(0xCB, 0xD5, 0xE1),
        align=PP_ALIGN.CENTER,
    )

    # 5 What AI pricing means
    slide = add_light_slide(prs)
    slide_title(slide, "What AI Pricing Means (and Does Not Mean)")
    accent_bar(slide)
    textbox(slide, Inches(0.7), Inches(1.55), Inches(5.5), Inches(0.4), "✓  We WILL do", size=20, bold=True, color=TEAL)
    bullets(
        slide,
        Inches(0.7),
        Inches(2.0),
        Inches(5.5),
        Inches(4.5),
        [
            "Rank products with highest profit upside.",
            "Propose safe price test bands (e.g. −5%, control, +5%).",
            "Run controlled experiments on live traffic.",
            "Learn over multiple test rounds toward optimum.",
            "Optimize for profit per visitor with floor rules.",
        ],
        size=15,
    )
    textbox(slide, Inches(6.8), Inches(1.55), Inches(5.5), Inches(0.4), "✗  We will NOT do (v1)", size=20, bold=True, color=RGBColor(0xDC, 0x26, 0x26))
    bullets(
        slide,
        Inches(6.8),
        Inches(2.0),
        Inches(5.5),
        Inches(4.5),
        [
            "Change catalog prices hourly without tests.",
            "Let LLM set prices with no statistical proof.",
            "Chase competitor prices automatically.",
            "Optimize conversion only (margin risk).",
        ],
        size=15,
    )

    # 6 How it works - 5 steps
    slide = add_light_slide(prs)
    slide_title(slide, "How It Works — 5 Steps")
    accent_bar(slide)
    steps = [
        ("1", "Analyze", "Sync Shopify orders + RipX analytics"),
        ("2", "Recommend", "AI ranks SKUs and suggests price bands"),
        ("3", "Approve", "Merchant reviews guardrails and launches"),
        ("4", "Test", "Controlled A/B on live traffic (checkout-safe)"),
        ("5", "Learn", "Pick winner → apply price → next round"),
    ]
    x_start = Inches(0.5)
    for i, (num, title, desc) in enumerate(steps):
        left = x_start + Inches(i * 2.45)
        circle = slide.shapes.add_shape(1, left, Inches(2.3), Inches(0.55), Inches(0.55))
        circle.fill.solid()
        circle.fill.fore_color.rgb = TEAL
        circle.line.fill.background()
        textbox(slide, left, Inches(2.38), Inches(0.55), Inches(0.4), num, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        textbox(slide, left - Inches(0.1), Inches(3.0), Inches(2.2), Inches(0.5), title, size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        textbox(slide, left - Inches(0.15), Inches(3.5), Inches(2.3), Inches(1.2), desc, size=13, color=SLATE, align=PP_ALIGN.CENTER)

    # 7 Architecture
    slide = add_light_slide(prs)
    slide_title(slide, "System Architecture — Four Layers")
    accent_bar(slide)
    layers = [
        ("Layer 1 — Intelligence", "Order data, SKU metrics, opportunity scoring"),
        ("Layer 2 — Experiment Design", "AI proposes SKUs, price arms, duration, guardrails"),
        ("Layer 3 — Execution", "Existing RipX price test + checkout resolver"),
        ("Layer 4 — Optimization Loop", "Profit analytics → winner → next test round"),
    ]
    y = Inches(1.7)
    for title, desc in layers:
        rect = slide.shapes.add_shape(1, Inches(0.7), y, Inches(11.8), Inches(1.05))
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        rect.line.color.rgb = TEAL
        textbox(slide, Inches(0.95), y + Inches(0.12), Inches(4), Inches(0.4), title, size=17, bold=True, color=NAVY)
        textbox(slide, Inches(5.2), y + Inches(0.18), Inches(7), Inches(0.5), desc, size=15, color=SLATE)
        y += Inches(1.25)

    # 8 Continuous optimization
    slide = add_light_slide(prs)
    slide_title(slide, "Continuous Optimization — How Learning Compounds")
    accent_bar(slide)
    textbox(slide, Inches(0.7), Inches(1.55), Inches(11.5), Inches(0.5), "Not one test forever — a learning loop:", size=18, bold=True, color=NAVY)
    bullets(
        slide,
        Inches(0.7),
        Inches(2.1),
        Inches(11.5),
        Inches(2.5),
        [
            "Round 1: Test control vs −5% vs +5%  →  winner: −5%",
            "Round 2: AI proposes −5%, −8%, −10% around the winner",
            "Round 3: Narrow band until profit stops improving",
            "Stop when guardrails hit or uplift is below threshold",
        ],
        size=17,
    )
    textbox(
        slide,
        Inches(0.7),
        Inches(4.8),
        Inches(11.5),
        Inches(1.2),
        "Later: bandit mode shifts traffic toward the best arm within a test — still within merchant-defined guardrails.",
        size=15,
        color=SLATE,
    )

    # 9 Example question
    slide = add_dark_slide(prs)
    slide_title(slide, "FAQ: Will it try 5%, then 10%, then 15%?", dark=True)
    accent_bar(slide, dark=True)
    textbox(
        slide,
        Inches(0.7),
        Inches(2.0),
        Inches(11.5),
        Inches(3.5),
        "Not all at once on day one.\n\n"
        "Each round tests a smart band together.\n"
        "After each winner, AI designs the NEXT round closer to the optimum.\n\n"
        "Fully hands-off auto-chaining is Phase 3 — with explicit merchant opt-in.",
        size=22,
        color=WHITE,
    )

    # 10 Why this approach
    slide = add_light_slide(prs)
    slide_title(slide, "Why This Approach Wins")
    accent_bar(slide)
    bullets(
        slide,
        Inches(0.7),
        Inches(1.6),
        Inches(11.5),
        Inches(5),
        [
            "Causal proof: real purchases, not model guesses.",
            "Merchant trust: approval gates + floor prices + max change limits.",
            "Profit-first: avoids discount-only bias of conversion metrics.",
            "Builds on RipX moat: checkout-accurate pricing (Discount Function path).",
            "Matches industry best practice (Shopify Smart Pricing, Intelligems, agency CRO).",
            "Echologyx DNA: experimentation execution at scale, productized.",
        ],
        size=17,
    )

    # 11 Already have
    slide = add_light_slide(prs)
    slide_title(slide, "What We Already Have")
    accent_bar(slide)
    bullets(
        slide,
        Inches(0.7),
        Inches(1.6),
        Inches(5.5),
        Inches(5),
        [
            "Price A/B tests (PDP + checkout)",
            "Discount Function + price resolver",
            "Revenue, profit, COGS metrics",
            "Sample size calculator",
            "Auto-stop on significance",
            "Shopify orders webhook + Admin API",
            "OpenAI integration (assistant layer)",
        ],
        size=16,
    )
    textbox(slide, Inches(6.8), Inches(2.0), Inches(5.5), Inches(3), "~80% of execution infrastructure is in place", size=22, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # 12 Need to build
    slide = add_light_slide(prs)
    slide_title(slide, "What We Need to Build")
    accent_bar(slide)
    bullets(
        slide,
        Inches(0.7),
        Inches(1.6),
        Inches(11.5),
        Inches(5),
        [
            "Order / SKU metrics store (historical sync + aggregation)",
            "Opportunity scoring engine (which products to optimize first)",
            "Price recommendation service (test bands + profit projection)",
            "Smart Pricing UI (primary flow; manual wizard → Advanced)",
            "Sequential learning orchestrator (next test after winner)",
            "Optional Phase 2+: Bayesian profit stats, bandit traffic allocation",
        ],
        size=17,
    )

    # 13 Roadmap
    slide = add_light_slide(prs)
    slide_title(slide, "Phased Roadmap")
    accent_bar(slide)
    phases = [
        ("Phase 0 — Foundation (2–3 wks)", "Order sync, SKU dashboard, read-only recommendations"),
        ("Phase 1 — MVP (4–6 wks)", "AI test proposal + one-click launch into RipX price tests"),
        ("Phase 2 — Learning loop (6–10 wks)", "Elasticity memory, profit Bayesian stats, Price Copilot"),
        ("Phase 3 — Continuous opt (10–16 wks)", "Auto next-round tests, bandit mode, winner apply workflow"),
    ]
    y = Inches(1.65)
    for title, desc in phases:
        textbox(slide, Inches(0.7), y, Inches(11.5), Inches(0.35), title, size=17, bold=True, color=TEAL)
        textbox(slide, Inches(0.9), y + Inches(0.38), Inches(11), Inches(0.35), desc, size=15, color=SLATE)
        y += Inches(1.15)

    # 14 Guardrails
    slide = add_light_slide(prs)
    slide_title(slide, "Guardrails & Trust")
    accent_bar(slide)
    bullets(
        slide,
        Inches(0.7),
        Inches(1.6),
        Inches(11.5),
        Inches(5),
        [
            "Minimum margin / floor price — never sell below COGS + threshold",
            "Maximum % price change per test cycle",
            "Promo / BFCM blackout windows",
            "Auto-pause if conversion or profit collapses",
            "Exposure cap (% of catalog revenue in active AI tests)",
            "Human approval before catalog price apply (v1)",
        ],
        size=17,
    )

    # 15 Tech stack
    slide = add_light_slide(prs)
    slide_title(slide, "Technology — No New Vendor Stack for MVP")
    accent_bar(slide)
    bullets(
        slide,
        Inches(0.7),
        Inches(1.6),
        Inches(5.5),
        Inches(4.5),
        [
            "Shopify Admin API (orders, products) — already integrated",
            "PostgreSQL + existing RipX analytics",
            "Bull/Redis job queue — already in stack",
            "OpenAI — optional for explanations (not pricing math)",
        ],
        size=16,
    )
    textbox(slide, Inches(6.8), Inches(1.6), Inches(5.5), Inches(0.4), "Optional later", size=18, bold=True, color=SLATE)
    bullets(
        slide,
        Inches(6.8),
        Inches(2.1),
        Inches(5.5),
        Inches(3.5),
        [
            "read_all_orders (deeper history)",
            "read_inventory (stock-aware pricing)",
            "Warehouse export (BigQuery)",
        ],
        size=15,
        color=SLATE,
    )

    # 16 Success metrics
    slide = add_light_slide(prs)
    slide_title(slide, "Success Metrics")
    accent_bar(slide)
    metrics = [
        ("Time to first AI price test", "< 5 minutes"),
        ("Tests launched via AI flow", "> 70% of new price tests"),
        ("Median profit lift on optimized SKUs", "+5% to +15%"),
        ("Guardrail false-positive pauses", "< 5%"),
    ]
    y = Inches(2.0)
    for label, target in metrics:
        textbox(slide, Inches(0.9), y, Inches(6), Inches(0.4), label, size=17, color=NAVY)
        textbox(slide, Inches(7.5), y, Inches(4.5), Inches(0.4), target, size=17, bold=True, color=TEAL, align=PP_ALIGN.RIGHT)
        y += Inches(0.85)

    # 17 Recommendation
    slide = add_dark_slide(prs)
    slide_title(slide, "Recommendation", dark=True)
    accent_bar(slide, dark=True)
    bullets(
        slide,
        Inches(0.7),
        Inches(1.7),
        Inches(11.5),
        Inches(4.5),
        [
            "Approve AI-first Smart Pricing as the primary RipX price product direction.",
            "Start with Phase 0 + Phase 1 MVP — validate with 1–2 pilot stores.",
            "Keep manual price tests as Advanced mode for power users.",
            "Position: “Profit-first AI pricing by the Echologyx experimentation team.”",
        ],
        size=20,
        color=WHITE,
    )

    # 18 Next steps
    slide = add_light_slide(prs)
    slide_title(slide, "Next Steps for Approval")
    accent_bar(slide)
    bullets(
        slide,
        Inches(0.7),
        Inches(1.6),
        Inches(11.5),
        Inches(5),
        [
            "Client sign-off on vision, phasing, and guardrail principles",
            "Confirm pilot store(s) and success criteria",
            "Kick off Phase 0: order sync + opportunity dashboard",
            "Deliver MVP demo in ~6 weeks for client review",
            "Plan GTM naming (e.g. ELX Smart Pricing / Splitter by Echologyx)",
        ],
        size=18,
    )

    # 19 Thank you
    slide = add_dark_slide(prs)
    textbox(slide, Inches(0.7), Inches(2.8), Inches(11.5), Inches(1), "Thank You", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(
        slide,
        Inches(0.7),
        Inches(4.0),
        Inches(11.5),
        Inches(0.8),
        "Questions & Discussion",
        size=24,
        color=TEAL,
        align=PP_ALIGN.CENTER,
    )
    textbox(
        slide,
        Inches(0.7),
        Inches(5.5),
        Inches(11.5),
        Inches(0.5),
        "Echologyx  •  RipX Smart Pricing Strategy  •  Confidential",
        size=14,
        color=RGBColor(0x94, 0xA3, 0xB8),
        align=PP_ALIGN.CENTER,
    )

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Wrote {OUTPUT}")


if __name__ == "__main__":
    build()
